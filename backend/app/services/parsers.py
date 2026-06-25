"""File content parsers — extract text from supported file types."""

import csv
import email
import io
import json
from pathlib import Path

import yaml

# Lazy imports for heavier libraries to speed up startup
_docx = None
_openpyxl = None
_pypdf = None
_pptx = None


def _get_docx():
    global _docx
    if _docx is None:
        import docx as _docx
    return _docx


def _get_openpyxl():
    global _openpyxl
    if _openpyxl is None:
        import openpyxl as _openpyxl
    return _openpyxl


def _get_pypdf():
    global _pypdf
    if _pypdf is None:
        import pypdf as _pypdf
    return _pypdf


def _get_pptx():
    global _pptx
    if _pptx is None:
        import pptx as _pptx
    return _pptx


# ─── Parse result ─────────────────────────────────────────────────────────────

class ParseResult:
    """Container for parsed text and metadata."""

    def __init__(self, text: str = "", warnings: list[str] | None = None, markers: list[dict] | None = None):
        self.text = text
        self.warnings = warnings or []
        self.markers = markers or []  # e.g. [{"type": "page", "label": "Page 1", "offset": 0}]


# ─── Dispatcher ───────────────────────────────────────────────────────────────

SUPPORTED_PARSE_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".yaml", ".yml",
    ".docx", ".xlsx", ".pdf", ".pptx", ".eml",
}


def parse_file(file_path: str) -> ParseResult:
    """Parse a file and return extracted text with metadata."""
    p = Path(file_path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = p.suffix.lower()
    parsers = {
        ".txt": _parse_text,
        ".md": _parse_text,
        ".csv": _parse_csv,
        ".json": _parse_json,
        ".yaml": _parse_yaml,
        ".yml": _parse_yaml,
        ".docx": _parse_docx,
        ".xlsx": _parse_xlsx,
        ".pdf": _parse_pdf,
        ".pptx": _parse_pptx,
        ".eml": _parse_eml,
    }

    parser = parsers.get(ext)
    if not parser:
        raise ValueError(f"Unsupported file type: {ext}")

    return parser(p)


# ─── Individual parsers ───────────────────────────────────────────────────────

def _parse_text(path: Path) -> ParseResult:
    """Parse plain text or markdown files."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return ParseResult("", warnings=[f"Read error: {e}"])
    return ParseResult(text)


def _parse_csv(path: Path) -> ParseResult:
    """Parse CSV into readable text."""
    warnings = []
    lines = []
    try:
        with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i == 0:
                    lines.append("Headers: " + " | ".join(row))
                elif i <= 500:
                    lines.append(" | ".join(row))
                else:
                    warnings.append(f"Truncated at 500 rows (file has more).")
                    break
    except Exception as e:
        return ParseResult("", warnings=[f"CSV parse error: {e}"])
    return ParseResult("\n".join(lines), warnings)


def _parse_json(path: Path) -> ParseResult:
    """Parse JSON into readable pretty-printed text."""
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        data = json.loads(raw)
        text = json.dumps(data, indent=2, ensure_ascii=False)
        if len(text) > 100_000:
            text = text[:100_000]
            return ParseResult(text, warnings=["JSON truncated at 100k characters."])
        return ParseResult(text)
    except json.JSONDecodeError:
        # Fall back to raw text
        return ParseResult(raw[:100_000], warnings=["Invalid JSON — stored as raw text."])
    except Exception as e:
        return ParseResult("", warnings=[f"JSON parse error: {e}"])


def _parse_yaml(path: Path) -> ParseResult:
    """Parse YAML into readable text."""
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        data = yaml.safe_load(raw)
        if data is None:
            return ParseResult(raw, warnings=["YAML loaded as empty — storing raw text."])
        text = yaml.dump(data, default_flow_style=False, allow_unicode=True)
        return ParseResult(text)
    except yaml.YAMLError:
        raw = path.read_text(encoding="utf-8", errors="replace")
        return ParseResult(raw[:100_000], warnings=["YAML parse failed — stored as raw text."])
    except Exception as e:
        return ParseResult("", warnings=[f"YAML error: {e}"])


def _parse_docx(path: Path) -> ParseResult:
    """Parse DOCX — extract paragraphs and tables."""
    docx = _get_docx()
    warnings = []
    parts = []
    try:
        doc = docx.Document(str(path))

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    parts.append(f"\n## {text}\n")
                else:
                    parts.append(text)

        for i, table in enumerate(doc.tables):
            parts.append(f"\n[Table {i + 1}]")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))

    except Exception as e:
        return ParseResult("", warnings=[f"DOCX parse error: {e}"])

    return ParseResult("\n".join(parts), warnings)


def _parse_xlsx(path: Path) -> ParseResult:
    """Parse XLSX — extract rows from all sheets."""
    openpyxl = _get_openpyxl()
    warnings = []
    parts = []
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"\n[Sheet: {sheet_name}]")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append(" | ".join(cells))
                    row_count += 1
                    if row_count >= 500:
                        warnings.append(f"Sheet '{sheet_name}' truncated at 500 rows.")
                        break
        wb.close()
    except Exception as e:
        return ParseResult("", warnings=[f"XLSX parse error: {e}"])

    return ParseResult("\n".join(parts), warnings)


def _parse_pdf(path: Path) -> ParseResult:
    """Parse PDF — extract page text."""
    pypdf = _get_pypdf()
    warnings = []
    parts = []
    markers = []
    try:
        reader = pypdf.PdfReader(str(path))
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            if page_text.strip():
                markers.append({"type": "page", "label": f"Page {i + 1}", "offset": len("\n".join(parts))})
                parts.append(f"[Page {i + 1}]")
                parts.append(page_text)
        if not parts:
            warnings.append("PDF text extraction returned empty — file may be image-based (OCR not supported).")
    except Exception as e:
        return ParseResult("", warnings=[f"PDF parse error: {e}"])

    return ParseResult("\n".join(parts), warnings, markers)


def _parse_pptx(path: Path) -> ParseResult:
    """Parse PPTX — extract slide text."""
    pptx = _get_pptx()
    warnings = []
    parts = []
    markers = []
    try:
        prs = pptx.Presentation(str(path))
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(cells))
            if slide_text:
                markers.append({"type": "slide", "label": f"Slide {i + 1}", "offset": len("\n".join(parts))})
                parts.append(f"[Slide {i + 1}]")
                parts.extend(slide_text)
    except Exception as e:
        return ParseResult("", warnings=[f"PPTX parse error: {e}"])

    return ParseResult("\n".join(parts), warnings, markers)


def _parse_eml(path: Path) -> ParseResult:
    """Parse EML email files."""
    warnings = []
    parts = []
    try:
        raw = path.read_bytes()
        msg = email.message_from_bytes(raw)

        # Headers
        for header in ["From", "To", "Subject", "Date"]:
            val = msg.get(header, "")
            if val:
                parts.append(f"{header}: {val}")

        parts.append("")

        # Body
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    charset = part.get_content_charset() or "utf-8"
                    payload = part.get_payload(decode=True)
                    if payload:
                        parts.append(payload.decode(charset, errors="replace"))
        else:
            charset = msg.get_content_charset() or "utf-8"
            payload = msg.get_payload(decode=True)
            if payload:
                parts.append(payload.decode(charset, errors="replace"))

    except Exception as e:
        return ParseResult("", warnings=[f"EML parse error: {e}"])

    return ParseResult("\n".join(parts), warnings)
